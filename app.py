# First, install the required packages by running this command in your terminal:
# pip install flask python-pptx requests pillow pymediainfo

from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches
import os
import tempfile
import requests
from PIL import Image
import io
from pymediainfo import MediaInfo

app = Flask(__name__)

def get_video_dimensions(video_path):
    """Get video dimensions using pymediainfo"""
    try:
        media_info = MediaInfo.parse(video_path)
        
        # Look for video tracks
        for track in media_info.tracks:
            if track.track_type == 'Video':
                width = track.width
                height = track.height
                
                if width and height:
                    return int(width), int(height)
        
        print("No video track found in the file")
        return None, None
        
    except Exception as e:
        print(f"Error getting video dimensions with pymediainfo: {e}")
        return None, None

@app.route('/')
def index():
    return "Hello, World! This is the PPT Video Injector API."

@app.route('/upload', methods=['POST'])
def upload_files():
    data = request.get_json()
    if not data or 'ppt' not in data or 'slides' not in data:
        return "Missing PPT URL or slides data", 400

    ppt_url = data['ppt']
    slides_info = data['slides']

    if not isinstance(slides_info, list) or not slides_info:
        return "Invalid slides data", 400

    # Create temporary directory
    with tempfile.TemporaryDirectory() as temp_dir:
        pptx_path = os.path.join(temp_dir, 'input.pptx')
        output_path = os.path.join(temp_dir, 'output.pptx')

        # Download PPTX from URL
        ppt_response = requests.get(ppt_url)
        if ppt_response.status_code != 200:
            return "Failed to download PPTX", 400
        with open(pptx_path, 'wb') as f:
            f.write(ppt_response.content)

        # Load the presentation
        prs = Presentation(pptx_path)
        num_slides = len(prs.slides)
        if num_slides == 0:
            return "Presentation has no slides", 400

        # Process each slide-video pair
        for idx, info in enumerate(slides_info):
            if not isinstance(info, dict) or 'number' not in info or 'videoLink' not in info:
                return f"Invalid slide info at index {idx}", 400

            slide_num = info['number']
            video_url = info['videoLink']

            if not isinstance(slide_num, int) or slide_num < 1 or slide_num > num_slides:
                return f"Invalid slide number {slide_num}", 400

            video_path = os.path.join(temp_dir, f'input_video_{slide_num}.mp4')
            thumbnail_path = os.path.join(temp_dir, f'thumbnail_{slide_num}.png')

            # Download video from URL
            video_response = requests.get(video_url)
            if video_response.status_code != 200:
                return f"Failed to download video for slide {slide_num}", 400
            with open(video_path, 'wb') as f:
                f.write(video_response.content)

            # Get actual video dimensions using pymediainfo
            video_width_px, video_height_px = get_video_dimensions(video_path)
            
            if video_width_px is None or video_height_px is None:
                # Fallback to default dimensions if we can't get actual dimensions
                print(f"Warning: Could not get video dimensions for slide {slide_num}, using defaults")
                video_width_px = 1920
                video_height_px = 1080
            
            print(f"Video dimensions: {video_width_px}x{video_height_px}")

            # Create a simple thumbnail based on video aspect ratio
            aspect_ratio = video_width_px / video_height_px
            if aspect_ratio > 1:  # Landscape
                thumb_width, thumb_height = 320, int(320 / aspect_ratio)
            else:  # Portrait or square
                thumb_width, thumb_height = int(240 * aspect_ratio), 240
            
            img = Image.new('RGB', (thumb_width, thumb_height), color='darkblue')
            img.save(thumbnail_path)

            # Convert pixel dimensions to inches (assuming 96 DPI)
            video_width_in = video_width_px / 96.0
            video_height_in = video_height_px / 96.0

            # Get slide dimensions in inches
            slide_width_in = prs.slide_width / 914400
            slide_height_in = prs.slide_height / 914400

            # Position 3: Bottom-right alignment (as in original code)
            left_in = slide_width_in - video_width_in
            top_in = slide_height_in - video_height_in

            # Ensure video doesn't exceed slide boundaries
            if video_width_in > slide_width_in or video_height_in > slide_height_in:
                # Scale down proportionally if video is larger than slide
                scale_factor = min(slide_width_in / video_width_in, slide_height_in / video_height_in) * 0.9  # 90% of slide max
                video_width_in *= scale_factor
                video_height_in *= scale_factor
                
                # Recalculate position after scaling (keep bottom-right alignment)
                left_in = slide_width_in - video_width_in
                top_in = slide_height_in - video_height_in
                print(f"Video scaled down by factor: {scale_factor:.2f}")

            # Ensure position is not negative
            left_in = max(0, left_in)
            top_in = max(0, top_in)

            print(f"Video size in inches: {video_width_in:.2f}x{video_height_in:.2f}")
            print(f"Position: left={left_in:.2f}, top={top_in:.2f}")

            # Add video to the specified slide
            slide = prs.slides[slide_num - 1]
            slide.shapes.add_movie(
                video_path,
                left=Inches(left_in),
                top=Inches(top_in),
                width=Inches(video_width_in),
                height=Inches(video_height_in),
                poster_frame_image=thumbnail_path,
                mime_type='video/mp4'
            )

        # Save the modified presentation
        prs.save(output_path)

        # Return the file
        return send_file(output_path, as_attachment=True, download_name='modified.pptx')


if __name__ == '__main__':
    app.run(debug=True)